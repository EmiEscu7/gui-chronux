import json
from abc import ABC, abstractmethod
from typing import Dict
import constants as ctes
import subprocess
import os
from Utils.loading import Loading
from pptx import Presentation
from pptx.util import Inches
from Utils.alert import Alert


class Analysis(ABC):

    def __init__(self, name):
        self._name = name
        self._parameters = None
        self._default_values = {}
        self._boxes = []
        self._files = None
        self._path_imgs = []
        self._presentation = None

    @property
    def files(self):
        return self._files

    @property
    def boxes(self):
        return self._boxes

    @property
    def default_values(self):
        return self._default_values

    @property
    def name(self):
        return self._name

    @property
    def parameters(self):
        return self._parameters

    @property
    def path_imgs(self):
        return self._path_imgs

    @default_values.setter
    def default_values(self, value):
        self._default_values = value

    @name.setter
    def name(self, value):
        self._name = value

    @parameters.setter
    def parameters(self, value):
        self._parameters = value

    @boxes.setter
    def boxes(self, value):
        self._boxes = value

    @files.setter
    def files(self, value):
        self._files = value

    @path_imgs.setter
    def path_imgs(self, value):
        self._path_imgs = value

    @abstractmethod
    def load_analysis(self) -> None:
        pass

    @abstractmethod
    def get_value_parameters(self) -> Dict:
        pass

    @abstractmethod
    def generate(self) -> None:
        pass

    @abstractmethod
    def save_params(self) -> None:
        pass

    @abstractmethod
    def destroy(self) -> None:
        pass

    @abstractmethod
    def generate_all_files(self, files):
        pass

    @abstractmethod
    def generate_all_files_th(self, files):
        pass

    def as_tuple(self, datos):
        return [(dato) for dato in datos]

    def show_params(self, master) -> None:
        self.boxes = self.parameters.load_params(master, [])
        for param in self.boxes:
            param.pack(padx=ctes.PADX_INPUTS, pady=ctes.PADY_INPUTS)

    def get_value_parameters(self) -> Dict:
        return self.parameters.get_data_params()

    def _save_signal(self, signal):
        file = f"{ctes.FOLDER_RES}Signal/signal.json"
        with open(file, 'w') as f:
            f.write(json.dumps({'signal': signal}))

    def _save_msignal(self, signal1, signal2):
        file = f"{ctes.FOLDER_RES}Signal/signal.json"
        with open(file, 'w') as f:
            f.write(json.dumps({'data1': signal1, 'data2': signal2}))

    def analysis(self, funtion, signal, params) -> float:
        self._save_signal(signal)
        process = subprocess.Popen(['matlab', '-batch', f"disp({funtion}({params}))"], stdout=subprocess.PIPE)

        output = process.communicate()[0]
        decode = output.decode()
        if "ERROR" in decode:
            print(decode)
            return 0.0
        try:
            result = float(decode.strip())
            return result
        except:
            print(decode.strip())
            Loading().change_state()
            return 0

    def manalysis(self, funtion, signal1, signal2, params) -> float:
        self._save_msignal(signal1, signal2)
        process = subprocess.Popen(['matlab', '-batch', f"disp({funtion}({params}))"], stdout=subprocess.PIPE)

        output = process.communicate()[0]
        decode = output.decode()
        if "ERROR" in decode:
            print(decode)
            return 0.0
        try:
            result = float(decode.strip())
            return result
        except:
            print(decode.strip())
            Loading().change_state()
            return 0

    @abstractmethod
    def generate_img_to_save(self, title, label):
        pass

    def generate_pptx(self, export_data_path):

        if self._presentation is None:
            self._presentation = Presentation()

        for label, path in self._path_imgs:

            layout = self._presentation.slide_layouts[5]
            slide = self._presentation.slides.add_slide(layout)

            slide.shapes.title.text = label

            # Add the plot image to the slide
            left = Inches(1.5)  # Adjust the positioning as needed
            top = Inches(1.5)
            height = Inches(5)
            slide.shapes.add_picture(path, left, top, height=height)
            os.remove(path)

        if self._presentation is not None:
            self._presentation.save(f'{export_data_path}/{self.files.info_file.file_name}.pptx')
            self._presentation = None
            Alert('Finished', f'File generated in {export_data_path}/{self.files.info_file.file_name}.pptx').show()
        self._path_imgs = []
