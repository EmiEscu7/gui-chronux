from Analysis.analysis import Analysis
from Parameters.spectogram_parameters import SpectogramParameters
from Files.file import File
from typing import List
import constants as ctes
import json
import numpy as np
import os
from Plots.Plot import Plot
from pptx import Presentation
from pptx.util import Inches
from Utils.loading import Loading

class SpectogramAnalysis(Analysis):

    def __init__(self):
        super().__init__(ctes.NAME_SPECTOGRAM)
        self._file_name = 'analysis'
        self._number_session = 0
        self._path_persist = './Persist/Parameters/spectogram_default'
        self._presentation = None
        self._export_data_path = './ExportData/Spectogram'
        self.data_compare = {}

    def _load_default_params(self, info_file) -> None:
        try:
            with open(f"{self._path_persist}-{info_file.file_name}.json", "r") as file:
                self.default_values = json.load(file)
        except:
            self.default_values = {
                'signal': info_file.signals[0],
                'movingwin1': '0.05',
                'movingwin2': '0.5',
                'taper1': '5',
                'taper2': '9',
                'sample_freq': '200',
                'freq': str(info_file.frequencies[0]),
                'freq_pass1': '0',
                'freq_pass2': '0',
                'time1': info_file.times[0],
                'time2': info_file.times[len(info_file.times) - 1],
                'trialave': '0',
                'err': '1',
            }

    def load_analysis(self) -> None:
        self._load_default_params(self.files.info_file)
        signals = (ctes.POPUP_MULTIPLE, ('Signal'), self.files.info_file.signals, self.default_values['signal'])
        check_all_signals = (ctes.CHECKBOX, 'All Signals', False, False)
        movingwin1 = (ctes.ENTRY, 'Moving Window 1', '', self.default_values['movingwin1'])
        movingwin2 = (ctes.ENTRY, 'Moving Window 2', '', self.default_values['movingwin2'])
        taper1 = (ctes.ENTRY, 'Taper 1', '', self.default_values['taper1'])
        taper2 = (ctes.ENTRY, 'Taper 2', '', self.default_values['taper2'])
        fs = (ctes.ENTRY, 'Frequency sample', '', self.default_values['sample_freq'])
        freqs = (ctes.POPUP, ('Frequency'), self.as_tuple(self.files.info_file.frequencies), self.default_values['freq'])
        freq_pass1 = (ctes.ENTRY, 'Frequency band 1', '', self.default_values['freq_pass1'])
        freq_pass2 =(ctes.ENTRY, 'Frequency band 2', '', self.default_values['freq_pass2'])
        time1 = (ctes.POPUP, ('Time 1'), self.files.info_file.times, self.default_values['time1'])
        time2 = (ctes.POPUP, ('Time 2'), self.files.info_file.times, self.default_values['time2'])
        trialave = (ctes.ENTRY, 'Trialave', '', self.default_values['trialave'])
        err = (ctes.ENTRY, 'Error', '', self.default_values['err'])

        self.parameters = SpectogramParameters(signals, check_all_signals, movingwin1, movingwin2, taper1, taper2, fs, freqs, freq_pass1, freq_pass2, time1, time2, trialave, err)


    def _generate_all(self, mw1, mw2, taper1, taper2, fs, freq, freq_pass1, freq_pass2, time1, time2, trialave, err) -> None:
        for signal, label in enumerate(self.files.info_file.signals):
            signal_matrix = self._get_signal_data(signal, freq, time1, time2, len(self.files.info_file.times))
            res = self.spectogram_analysis(signal_matrix, mw1, mw2, taper1, taper2, fs, freq_pass1, freq_pass2, trialave, err)
            if res == 1:
                self._generate_pptx(f"{label} - Spectogram Plot", label)

            if self._presentation is not None:
                self._presentation.save(f'{self._export_data_path}/{self.files.info_file.file_name}.pptx')
                self._presentation = None

    def generate_all_files(self, files):
        Loading().start(self.generate_all_files_th, (files,))

    def generate_all_files_th(self, files):
        for file in files:
            current_file = self.files.get_specific_file(file)
            if current_file.get_parameters() is not None and len(current_file.get_parameters()) > 0:
                data = current_file.get_parameters()
            else:
                data = self.get_value_parameters()
            movingwin1 = data['movingwin1']
            movingwin2 = data['movingwin2']
            taper1 = data['taper1']
            taper2 = data['taper2']
            fs = data['fs']
            str_freq = data['freq']
            freq = self.files.info_file.frequencies.index(str_freq) + 1
            freq_pass1 = data['freq_pass1']
            freq_pass2 = data['freq_pass2']
            str_time1 = data['time1']
            time1 = self.files.info_file.times.index(str_time1)
            str_time2 = data['time2']
            time2 = self.files.info_file.times.index(str_time2)
            trialave = data['trialave']
            err = data['err']
            str_signal = data['signal']
            signal = self.files.info_file.signals.index(str_signal)
            signal_matrix = self._get_signal_data(signal, freq, time1, time2, len(self.files.info_file.times))
            res = self.spectogram_analysis(signal_matrix, movingwin1, movingwin2, taper1, taper2, fs, freq_pass1, freq_pass2, trialave, err)
            if res == 1:
                self._save_data_temp(file)

        self._generate_plot_all_files(current_file.info_file.file_name)

    def _save_data_temp(self, name_file):
        file = f"{ctes.FOLDER_RES}Spectogram/{self._file_name}.json"
        # Open file in read mode
        with open(file, 'r') as f:
            # read file content
            content = f.read()

            # decdoe content ot json format
            data = json.loads(content)

        os.remove(file)
        to_save = {
            't': data['t'],
            'f': data['f'],
            'S': 10 * np.log10(data['S'])
        }
        self.data_compare[name_file] = to_save

    def _generate_plot_all_files(self, title):
        self._number_session += 1
        Plot().add_multi_color_plot(self.data_compare, 'Time (s)', 'Frequency (Hz)', 'Signal Spectogram', title)
        Loading().change_state()

    def generate(self) -> None:
        Loading().start(self.generate_th)

    def generate_th(self):
        data = self.get_value_parameters()
        movingwin1 = data['movingwin1']
        movingwin2 = data['movingwin2']
        taper1 = data['taper1']
        taper2 = data['taper2']
        fs = data['fs']
        str_freq = data['freq']
        freq = self.files.info_file.frequencies.index(str_freq) + 1
        freq_pass1 = data['freq_pass1']
        freq_pass2 = data['freq_pass2']
        str_time1 = data['time1']
        time1 = self.files.info_file.times.index(str_time1)
        str_time2 = data['time2']
        time2 = self.files.info_file.times.index(str_time2)
        trialave = data['trialave']
        err = data['err']
        all_signals = data['all']
        if all_signals:
            self._generate_all(movingwin1, movingwin2, taper1, taper2, fs, freq, freq_pass1, freq_pass2, time1, time2,trialave, err)
        else:
            str_signal = str(data['signal'])
            arr_signal = str_signal.split(",")
            if len(arr_signal) == 1:
                signal = self.files.info_file.signals.index(str_signal)
                signal_matrix = self._get_signal_data(signal, freq, time1, time2, len(self.files.info_file.times))
                res = self.spectogram_analysis(signal_matrix, movingwin1, movingwin2, taper1, taper2, fs, freq_pass1,freq_pass2, trialave, err)
                if res == 1:
                    self._generate_plot(f"{self._number_session} - Spectogram Plot")
                Loading().change_state()
            else:
                for select_signal in arr_signal:
                    signal = self.files.info_file.signals.index(select_signal.strip())
                    signal_matrix = self._get_signal_data(signal, freq, time1, time2, len(self.files.info_file.times))
                    res = self.spectogram_analysis(signal_matrix, movingwin1, movingwin2, taper1, taper2, fs, freq_pass1, freq_pass2, trialave, err)
                    if res == 1:
                        self._save_data_temp(select_signal)
                self._generate_plot_all_files(self.files.info_file.file_name)



    def _get_signal_data(self, signal, freq, time1, time2, n, file = None) -> str:
        if file is None:
            data_in_freq = self.files.info_file.nex
        else:
            data_in_freq = file.nex
        range1 = self._get_range(signal, time1, n) - 1
        range2 = self._get_range(signal, time2, n)
        columns = data_in_freq.iloc[range1:range2]
        if freq is None:
            data = columns.iloc[:]
        else:
            data = columns.iloc[:freq]
        matlab_string = "["
        for r, fila in data.iterrows():
            matlab_string += " ".join(map(str, fila)) + "; "
        matlab_string = matlab_string[:-2]  # Eliminar el último "; "
        matlab_string += "]"
        return matlab_string


    def _get_range(self, i, car, n) -> int:
        return (2 + n * i) + car

    def spectogram_analysis(self, signal, movingwin1, movingwin2, taper1, taper2, fs, freq_pass1, freq_pass2, trialave, err):
        # Execute MATLAB in CMD and capture output
        movingwin = f'[{movingwin1} {movingwin2}]'
        tapers = f'[{taper1} {taper2}]'
        fpass = f'[{freq_pass1} {freq_pass2}]'
        params = f"{movingwin}, {tapers}, {fpass}, {fs}, {err}, {trialave}, '{ctes.FOLDER_RES + 'Spectogram/'}', '{self._file_name}'"
        # function = f"SpectogramAnalysis({signal}, {movingwin}, {tapers}, {fpass}, {fs}, {err}, {trialave}, '{ctes.FOLDER_RES + 'Spectogram/'}', '{self._file_name}')"
        return self.analysis('SpectogramAnalysis', signal, params)

    def _generate_plot(self, title):
        file = f"{ctes.FOLDER_RES}Spectogram/{self._file_name}.json"
        # Open file in read mode
        with open(file, 'r') as f:
            # read file content
            content = f.read()

            # decdoe content ot json format
            data = json.loads(content)

        os.remove(file)

        # get data
        s = data['S']
        t = data['t']
        f = data['f']

        # show plot
        self._number_session += 1
        Plot().add_color_plot(t, f, 10 * np.log10(s), 'Time (s)', 'Frequency (Hz)', 'Signal Spectogram', title)

    def save_params_session(self) -> None:
        data = self.get_value_parameters()
        signal = data['signal']
        movingwin1 = data['movingwin1']
        movingwin2 = data['movingwin2']
        taper1 = data['taper1']
        taper2 = data['taper2']
        fs = data['fs']
        freq = data['freq']
        freq_pass1 = data['freq_pass1']
        freq_pass2 = data['freq_pass2']
        time1 = data['time1']
        time2 = data['time2']
        trialave = data['trialave']
        err = data['err']

        self.default_values = {
            'signal': signal,
            'movingwin1': str(movingwin1),
            'movingwin2': str(movingwin2),
            'taper1': str(taper1),
            'taper2': str(taper2),
            'sample_freq': str(fs),
            'freq': str(freq),
            'freq_pass1': str(freq_pass1),
            'freq_pass2': str(freq_pass2),
            'time1': str(time1),
            'time2': str(time2),
            'trialave': str(trialave),
            'err': str(err),
        }

    def save_params(self) -> None:
        self.save_params_session()

        with open(f"{self._path_persist}-{self.files.info_file.file_name}.json", "w") as file:
            json.dump(self.default_values, file)

    def destroy(self) -> None:
        self.parameters.destroy()
        for box in self.boxes:
            box.destroy()

    def _generate_pptx(self, title, label):
        if self._presentation is None:
            self._presentation = Presentation()

        file = f"{ctes.FOLDER_RES}Spectogram/{self._file_name}.json"
        with open(file, 'r') as f:
            # read file content
            content = f.read()

            # decdoe content ot json format
            data = json.loads(content)

        os.remove(file)

        # get data
        s = data['S']
        t = data['t']
        f = data['f']
        path_img = f"{self._export_data_path}/{label}.png"
        Plot().get_color_plot(t,f, 10 * np.log10(s),  'Time (s)', 'Frequency (Hz)', 'Signal Spectogram',
                        path_img)

        layout = self._presentation.slide_layouts[5]
        slide = self._presentation.slides.add_slide(layout)

        slide.shapes.title.text = label

        # Add the plot image to the slide
        left = Inches(1.5)  # Adjust the positioning as needed
        top = Inches(1.5)
        height = Inches(5)
        slide.shapes.add_picture(path_img, left, top, height=height)
        os.remove(path_img)
